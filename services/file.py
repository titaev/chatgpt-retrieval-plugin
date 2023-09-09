import os
from io import BufferedReader
from typing import Optional
from fastapi import UploadFile
import mimetypes
from PyPDF2 import PdfReader
import docx2txt
import csv
import pptx
import textract
import tempfile
import chardet

from models.models import Document, DocumentMetadata, Source


async def get_document_from_file(file: UploadFile) -> Document:
    extracted_text = await extract_text_from_form_file(file)
    metadata = DocumentMetadata(
        source=Source.file,
    )
    doc = Document(text=extracted_text, metadata=metadata)

    return doc


def extract_text_from_filepath(filepath: str, mimetype: Optional[str] = None) -> str:
    """Return the text content of a file given its filepath."""

    if mimetype is None:
        # Get the mimetype of the file based on its extension
        mimetype, _ = mimetypes.guess_type(filepath)

    if not mimetype:
        if filepath.endswith(".md"):
            mimetype = "text/markdown"
        else:
            raise Exception("Unsupported file type")

    # Open the file in binary mode
    file = open(filepath, "rb")
    extracted_text = extract_text_from_file(file, mimetype)

    return extracted_text


def extract_text_from_file(file: BufferedReader, mimetype: str) -> str:
    if mimetype == "application/pdf":
        # Extract text from pdf using PyPDF2
        reader = PdfReader(file)
        extracted_text = " ".join([page.extract_text() for page in reader.pages])
    elif mimetype == "text/plain" or mimetype == "text/plain;charset=utf-8" or mimetype == "text/markdown":
        # Read text from plain text file
        extracted_text = file.read().decode("utf-8")
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    ):
        # Extract text from docx using docx2txt
        extracted_text = docx2txt.process(file)
    elif mimetype == "application/msword":
        # Extract text from doc using textract
        with tempfile.NamedTemporaryFile(delete=False) as temp:
            temp.write(file.read())
        try:
            raw_text = textract.process(temp.name, extension='doc')
            # Detect encoding
            detected = chardet.detect(raw_text)
            encoding = detected.get("encoding", "utf-8")
            # Decode the text using the detected encoding
            extracted_text = raw_text.decode(encoding)
        finally:
            os.unlink(temp.name)
    elif mimetype == "text/csv":
        # Extract text from csv using csv module
        extracted_text = ""
        decoded_buffer = (line.decode("utf-8") for line in file)
        reader = csv.reader(decoded_buffer)
        for row in reader:
            extracted_text += " ".join(row) + "\n"
    elif (
        mimetype
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    ):
        # Extract text from pptx using python-pptx
        extracted_text = ""
        presentation = pptx.Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            extracted_text += run.text + " "
                    extracted_text += "\n"
    else:
        # Unsupported file type
        file.close()
        raise ValueError("Unsupported file type: {}".format(mimetype))

    file.close()
    return extracted_text


def save_temp_file(file_stream):
    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)
    return temp_file_path


# Extract text from a file based on its mimetype
async def extract_text_from_form_file(file: UploadFile):
    """Return the text content of a file."""
    # get the file body from the upload file object
    mimetype = file.content_type
    print(f"mimetype: {mimetype}")
    print(f"file.file: {file.file}")
    print("file: ", file)

    file_stream = await file.read()

    temp_file_path = "/tmp/temp_file"

    # write the file to a temporary location
    with open(temp_file_path, "wb") as f:
        f.write(file_stream)

    try:
        extracted_text = extract_text_from_filepath(temp_file_path, mimetype)
    except Exception as e:
        print(f"Error: {e}")
        os.remove(temp_file_path)
        raise e

    # remove file from temp location
    os.remove(temp_file_path)

    return extracted_text

#
# class TableFilesSupport:
#     """support csv or excel"""
#
#     excel_mimetype = ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "application/vnd.ms-excel")
#     csv_mimetype = ("text/csv", )
#
#     def is_table_file(self):
#         return True
#
#     async def get_documents_from_table_file(self, file: UploadFile, mimetype: str):
#         """support csv or excel"""
#         file_stream = await file.read()
#
#         # START of sync part, should be sync because of temp_file_path same for all
#         temp_file_path = save_temp_file(file_stream)
#         if mimetype in self.excel_mimetype:
#             return self.get_documents_from_excel(temp_file_path)
#
#         elif mimetype in self.csv_mimetype:
#             return self.get_documents_from_csv(temp_file_path)
#         # END of sync part
#
#         raise
#
#     def get_documents_from_excel(self, temp_file_path):
#         loader = UnstructuredExcelLoader(temp_file_path, mode="elements")
#         docs = loader.load()
#         return docs
#
#     def get_documents_from_csv(self, temp_file_path):
#         loader = CSVLoader(file_path=temp_file_path)
#         docs = loader.load()
#         return docs

